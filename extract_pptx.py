from pptx import Presentation
import sys
import os

def extract_text_and_tables(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return
        
    prs = Presentation(pptx_path)
    output = []
    
    for slide_idx, slide in enumerate(prs.slides):
        output.append(f"\n## Slide {slide_idx + 1}")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Extract text
                output.append(shape.text.strip())
                
            if shape.has_table:
                # Extract table
                table = shape.table
                table_md = []
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for cell in row.cells:
                        # Clean cell text (remove newlines inside cells)
                        cell_text = cell.text_frame.text.replace("\n", " ").replace("\r", " ").strip()
                        row_data.append(cell_text)
                    
                    row_str = "| " + " | ".join(row_data) + " |"
                    table_md.append(row_str)
                    
                    # Add separator after header
                    if row_idx == 0:
                        separator = "| " + " | ".join(["---"] * len(row_data)) + " |"
                        table_md.append(separator)
                        
                output.append("\n".join(table_md))
                
    return "\n\n".join(output)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <file.pptx>")
        sys.exit(1)
        
    md_content = extract_text_and_tables(sys.argv[1])
    
    out_file = sys.argv[1].replace(".pptx", "_extracted.md")
    with open(out_file, "w", encoding="utf-8") as f:
        f.write(md_content)
        
    print(f"Extracted to {out_file}")
