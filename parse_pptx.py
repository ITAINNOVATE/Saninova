from pptx import Presentation

prs = Presentation("QUANTIFICATION DES BESOINS.pptx")
for i, slide in enumerate(prs.slides):
    title = slide.shapes.title.text if slide.shapes.title else "NO TITLE"
    print(f"--- Slide {i+1} ---")
    print(f"TITLE: {title}")
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            print(f"TEXT: {shape.text[:100]}...")
    if i > 10:
        break
